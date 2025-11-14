from fastapi import FastAPI, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import List, Optional
from fastapi import Query
import os
import re
import shutil
import tempfile
import fitz
import docx
import pptx
from langchain_groq import ChatGroq
from langchain.prompts import PromptTemplate
from fpdf import FPDF
import uuid
import asyncio
import datetime

# Set your Groq API key
os.environ['GROQ_API_KEY'] = "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"

# Initialize FastAPI app
app = FastAPI()

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # In production, replace with your frontend domain
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Instead of using a temporary directory
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "outputs")
os.makedirs(OUTPUT_DIR, exist_ok=True)

# === File Handlers === #
def extract_text_from_pdf(file_path):
    text = ""
    doc = fitz.open(file_path)
    for page in doc:
        text += page.get_text()
    doc.close()
    return text

def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_text_from_docx(file_path):
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = pptx.Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def save_to_pdf(text, output_path):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("helvetica", size=10)
    for line in text.split('\n'):
        if line.startswith('###'):
            pdf.set_font("helvetica", 'B', 14)
            pdf.cell(0, 10, line.replace('###', '').strip(), ln=True)
            pdf.set_font("helvetica", size=10)
        elif line.strip():
            pdf.multi_cell(0, 5, line)
    pdf.output(output_path)
    return output_path

# === LangChain Setup === #
llm = ChatGroq(model='llama-3.3-70b-versatile')

def get_mcq_prompt_template(difficulty: str = "medium", solution_variant: int = 1):
    difficulty_instructions = {
        "easy": "Generate simple, direct questions with clear correct answers based only on facts explicitly stated in the text.",
        "medium": "Generate moderately complex questions that require interpretation, context understanding, or applying information from the text.",
        "hard": "Generate complex, analytical questions that require deep understanding, reasoning, or synthesis of ideas from multiple parts of the text."
    }

    variation_instructions = {
        1: "Focus on factual recall and basic comprehension questions.",
        2: "Focus on application and analysis questions.",
        3: "Focus on synthesis and evaluation questions."
    }

    return PromptTemplate.from_template(f"""
You are a professional educator creating multiple distinct solution sets for the same material.

For Solution Set {solution_variant}, generate {{n}} Multiple Choice Questions (MCQs) with these requirements:

- Each question must be UNIQUE and not repeated from other solution sets
- Questions should focus on different aspects of the text than other solution sets
- Difficulty Level: {difficulty.upper()} - {difficulty_instructions[difficulty]}
- Specific Focus: {variation_instructions.get(solution_variant, variation_instructions[1])}
- Each question must have exactly 4 answer options: A, B, C, D
- Only ONE correct option per question
- Options must be in DIFFERENT ORDER than other solution sets
- Questions must be WORDED DIFFERENTLY than in other solution sets
- Correct answers may test different aspects of the same concept

Format each question as:

### Question [number]
[Unique question text different from other solution sets]
A) [Option A - different wording/order]
B) [Option B - different wording/order]
C) [Option C - different wording/order]
D) [Option D - different wording/order]
**Answer: [Different letter than other solution sets when testing same concept]**

### Document Extract:
{{document_text}}
""")

def get_descriptive_prompt_template(difficulty: str = "medium", solution_variant: int = 1):
    difficulty_instructions = {
        "easy": "Ask basic questions that can be directly answered using simple sentences from the text.",
        "medium": "Ask thoughtful questions that require connecting multiple ideas from the text.",
        "hard": "Ask challenging questions that require reasoning, inference, or critical analysis based on the text."
    }

    variation_instructions = {
        1: "Focus on factual and definition questions.",
        2: "Focus on process and relationship questions.",
        3: "Focus on implication and evaluation questions."
    }

    return PromptTemplate.from_template(f"""
You are a professional educator creating multiple distinct solution sets for the same material.

For Solution Set {solution_variant}, generate {{n}} short descriptive questions with these requirements:

- Each question must be UNIQUE and not repeated from other solution sets
- Questions should approach the material from different angles than other solution sets
- Difficulty Level: {difficulty.upper()} - {difficulty_instructions[difficulty]}
- Specific Focus: {variation_instructions.get(solution_variant, variation_instructions[1])}
- Questions must be WORDED DIFFERENTLY than in other solution sets
- Answers should demonstrate different perspectives on the same concepts

Format each question as:

### Question [number]
[Unique question text different from other solution sets]

### Answer
[Comprehensive answer from a different perspective than other solution sets]

### Document Extract:
{{document_text}}
""")


def get_assignment_prompt_template():
    return PromptTemplate.from_template("""
    You are an expert at analyzing text and generating useful study materials. For the given document extract below, perform the following task:

    Generate {n} Assignment Questions:
    - Create {n} assignment questions that require deeper understanding and application of concepts from the text.
    - Include a mix of problem-solving, analysis, and critical thinking questions.
    - Format as a numbered list.
    - Include a detailed answer for each question that demonstrates what an ideal response would look like.
    
    ### Document Extract:
    {document_text}
    """)

def get_summary_prompt_template():
    return PromptTemplate.from_template("""
    You are an expert at analyzing text and generating useful study materials. For the given document extract below, perform the following task:

    Write a Comprehensive Summary:
    - Summarize the document extract in 2-3 paragraphs, capturing the main ideas, key points, and purpose.
    - Include a concise executive summary at the beginning (1-2 paragraphs).
    - Include a section on key concepts or terminology introduced in the document.
    - Conclude with the main takeaways from the document.
    
    ### Document Extract:
    {document_text}
    """)

# === API Models === #
class GenerateRequest(BaseModel):
    title: Optional[str] = None
    numberOfQuestions: int = 10
    type: str  # "mcq", "descriptive", "assignment", "summary"
    difficulty: str = "medium"  # "easy", "medium", "hard"
    numberOfSolutions: int = 1  # 1-3



# === API Endpoints === #
@app.post("/api/upload/")
async def upload_files(files: List[UploadFile] = File(...)):
    file_ids = []
    
    for file in files:
        file_id = str(uuid.uuid4())
        file_path = os.path.join(OUTPUT_DIR, f"{file_id}_{file.filename}")
        
        try:
            with open(file_path, "wb") as f:
                shutil.copyfileobj(file.file, f)
            file_ids.append({"id": file_id, "name": file.filename, "path": file_path})
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"Error uploading file: {e}")
    
    return {"files": file_ids}




@app.post("/api/generate/")
async def generate_content(
    request: GenerateRequest,
    file_ids: str = Query(None)
):
    try:
        # Convert string to list if provided
        file_id_list = file_ids.split(',') if file_ids else []

        if not file_id_list:
            raise HTTPException(status_code=400, detail="No file IDs provided")

        # Validate file IDs
        file_paths = []
        for file_id in file_id_list:
            matching_files = [f for f in os.listdir(OUTPUT_DIR) if f.startswith(file_id)]
            if not matching_files:
                raise HTTPException(status_code=404, detail=f"File with ID {file_id} not found")
            file_paths.append(os.path.join(OUTPUT_DIR, matching_files[0]))

        # Extract text from files
        full_text = ""
        for path in file_paths:
            ext = path.split('.')[-1].lower()
            if ext == "pdf":
                full_text += extract_text_from_pdf(path)
            elif ext == "txt":
                full_text += extract_text_from_txt(path)
            elif ext in ["doc", "docx"]:
                full_text += extract_text_from_docx(path)
            elif ext in ["ppt", "pptx"]:
                full_text += extract_text_from_pptx(path)
            full_text += "\n\n"

        # Generate content based on type
        result_content = ""

        if request.type in ["mcq", "descriptive"]:
            prompt_func = get_mcq_prompt_template if request.type == "mcq" else get_descriptive_prompt_template
            combined_solutions = []
            for i in range(request.numberOfSolutions):
                prompt = prompt_func(request.difficulty)
                chain = prompt | llm

                # Add variation prompt to encourage unique solutions
                result = await asyncio.to_thread(chain.invoke, {
                    "n": request.numberOfQuestions,
                    "document_text": f"{full_text}\n\n(Variant {i + 1})"
                })

                combined_solutions.append(f"### Solution {i + 1}\n{result.content.strip()}")

            result_content = "\n\n".join(combined_solutions)

        elif request.type == "assignment":
            prompt = get_assignment_prompt_template()
            chain = prompt | llm
            result = await asyncio.to_thread(chain.invoke, {
                "n": request.numberOfQuestions,
                "document_text": full_text
            })
            result_content = result.content

        elif request.type == "summary":
            prompt = get_summary_prompt_template()
            chain = prompt | llm
            result = await asyncio.to_thread(chain.invoke, {
                "document_text": full_text
            })
            result_content = result.content

        else:
            raise HTTPException(status_code=400, detail="Invalid content type")

        # Generate a timestamp
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")

        # Create a title for the output file
        title = request.title if request.title else f"{request.type}_{timestamp}"

        # Save as PDF
        pdf_name = f"{title.replace(' ', '_')}_{timestamp}.pdf"
        pdf_path = os.path.join(OUTPUT_DIR, pdf_name)
        save_to_pdf(result_content, pdf_path)

        # Return result
        return {
            "content": result_content,
            "pdf_url": f"/api/download/{pdf_name}",
            "title": title,
            "quizId": str(uuid.uuid4()) if request.type == "quiz" else None
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/api/download/{filename}")
async def download_file(filename: str):
    file_path = os.path.join(OUTPUT_DIR, filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail=f"File not found at {file_path}")
    return FileResponse(file_path, filename=filename)

@app.get("/")
async def root():
    return {"message": "Study AI API is running"}

# Run with: uvicorn api:app --reload
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)